"""
Цветные презентации PPTX: иллюстрации (нейро + веб), заметки докладчика.
Сборка совместима с PowerPoint и Keynote (OOXML, безопасный текст, PNG/JPEG для картинок).
"""

from __future__ import annotations

import asyncio
import json
import logging
import re
import uuid
import zipfile
from pathlib import Path
from typing import Any, Optional

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.config import settings
from app.image_utils import image_api_response_to_sse_href
from app.mws_client import MWSClient
from app.web_tools import image_search_ddg_urls

logger = logging.getLogger("gpthub.presentation")

# Символы, недопустимые в XML 1.0 / OOXML (Keynote иначе может отказать открыть файл).
_OOXML_ILLEGAL = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")
_CYRILLIC_RE = re.compile(r"[\u0400-\u04FF]")
_NEURO_SLIDE_JUNK = re.compile(
    r"\b(presentation|slideshow|slide\s*deck|powerpoint|keynote|infographic|"
    r"bullet\s*points?|title\s*bar|layout\s+with|chart\s+with\s+text|"
    r"with\s+(?:a\s+)?title|mock\s*ups?|ui\s+elements?|poster\s+with\s+text)\b",
    re.I,
)


def sanitize_ooxml_text(s: str, max_len: int = 32000) -> str:
    if not s:
        return ""
    t = _OOXML_ILLEGAL.sub("", str(s))
    t = "".join(c for c in t if not (0xD800 <= ord(c) <= 0xDFFF))
    return t[:max_len]


# Белый список имён шрифтов для полей title_font / body_font в JSON слайда (Keynote/PowerPoint).
_SAFE_FONT_ALIASES: dict[str, str] = {
    "arial": "Arial",
    "calibri": "Calibri",
    "cambria": "Cambria",
    "georgia": "Georgia",
    "helvetica": "Helvetica",
    "tahoma": "Tahoma",
    "times": "Times New Roman",
    "times new roman": "Times New Roman",
    "verdana": "Verdana",
}


def _font_face_from_field(raw: Any, fallback: str) -> str:
    s = str(raw or "").strip().lower()
    if not s:
        return fallback
    return _SAFE_FONT_ALIASES.get(s) or _SAFE_FONT_ALIASES.get(s.replace("  ", " ")) or fallback


def _bundled_template_path() -> Path:
    return Path(__file__).resolve().parent / "assets" / "keynote_base.pptx"


def _open_presentation_base() -> Presentation:
    """Каркас: по умолчанию голый Presentation() — проще для Keynote; иначе файл шаблона."""
    if not settings.gena_pptx_use_bundled_template and not (settings.gena_pptx_template_path or "").strip():
        return Presentation()
    custom = (settings.gena_pptx_template_path or "").strip()
    if custom:
        p = Path(custom)
        if p.is_file():
            try:
                return Presentation(str(p))
            except Exception as e:
                logger.warning("custom PPTX template unreadable, fallback: %s (%s)", p, e)
    bundled = _bundled_template_path()
    if bundled.is_file():
        try:
            return Presentation(str(bundled))
        except Exception as e:
            logger.warning("bundled PPTX template unreadable, fallback empty deck: %s (%s)", bundled, e)
    return Presentation()


def _validate_ooxml_package(path: Path, slide_count: int) -> bool:
    """Минимальная проверка ZIP-пакета презентации."""
    required = (
        "[Content_Types].xml",
        "_rels/.rels",
        "ppt/presentation.xml",
    )
    try:
        with zipfile.ZipFile(path, "r") as z:
            names = set(z.namelist())
            for req in required:
                if req not in names:
                    logger.error("pptx package missing %s", req)
                    return False
            for i in range(1, max(1, slide_count) + 1):
                sp = f"ppt/slides/slide{i}.xml"
                if sp not in names:
                    logger.error("pptx package missing %s", sp)
                    return False
    except zipfile.BadZipFile as e:
        logger.error("pptx not a valid zip: %s", e)
        return False
    return True


def _roundtrip_repair(path: Path) -> None:
    """Перечитать и сохранить — часто устраняет мелкие несовместимости импортёров."""
    prs = Presentation(str(path))
    prs.save(str(path))


def _ensure_keynote_safe_image(src: Optional[Path]) -> Optional[Path]:
    """Всегда пересохраняем в PNG RGB — Keynote/PowerPoint меньше ломаются на «чужих» JPEG/WebP."""
    if src is None or not src.is_file():
        return None
    max_px = int(settings.gena_pptx_max_image_px)
    try:
        from PIL import Image

        with Image.open(src) as im0:
            im0.verify()

        with Image.open(src) as im:
            im.load()
            im = im.convert("RGB")
            w, h = im.size
            if max(w, h) > max_px:
                scale = max_px / float(max(w, h))
                w, h = max(1, int(w * scale)), max(1, int(h * scale))
                try:
                    resample = Image.Resampling.LANCZOS
                except AttributeError:
                    resample = Image.LANCZOS  # Pillow < 9.1
                im = im.resize((w, h), resample)
            out = src.parent / f"{src.stem}_embed_{uuid.uuid4().hex[:8]}.png"
            im.save(out, "PNG", optimize=True)
            return out
    except Exception as e:
        logger.warning("cannot use image for PPTX (convert/skip): %s (%s)", src, e)
        return None


def _effective_preset(slide_data: dict[str, Any], idx: int) -> dict[str, Any]:
    """Пресет + font_scale и переопределения шрифтов из JSON слайда."""
    preset = dict(_preset_for_slide(slide_data, idx))
    raw_scale = slide_data.get("font_scale")
    if raw_scale is None:
        raw_scale = slide_data.get("typography_scale")
    try:
        sc = float(raw_scale)
    except (TypeError, ValueError):
        sc = 1.0
    sc = max(0.75, min(1.35, sc))
    for k in ("title_pt", "subtitle_pt", "body_pt", "footer_pt"):
        preset[k] = max(8, min(56, int(round(int(preset[k]) * sc))))
    preset["title_face"] = _font_face_from_field(slide_data.get("title_font"), preset["title_face"])
    preset["body_face"] = _font_face_from_field(slide_data.get("body_font"), preset["body_face"])
    preset["notes_face"] = _font_face_from_field(slide_data.get("notes_font"), preset["notes_face"])
    return preset


_DEFAULT_ACCENTS = [
    "#1e40af",
    "#6d28d9",
    "#0f766e",
    "#c2410c",
    "#be185d",
    "#15803d",
    "#b45309",
    "#0369a1",
]


def _hex_to_rgb(h: str | None, idx: int) -> tuple[RGBColor, RGBColor]:
    """Акцент + светлый фон под контент."""
    raw = (h or "").strip().lstrip("#") or _DEFAULT_ACCENTS[idx % len(_DEFAULT_ACCENTS)].lstrip("#")
    if len(raw) != 6 or not re.fullmatch(r"[0-9a-fA-F]{6}", raw):
        raw = _DEFAULT_ACCENTS[idx % len(_DEFAULT_ACCENTS)].lstrip("#")
    r, g, b = int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16)
    accent = RGBColor(r, g, b)
    br = min(255, r // 5 + 204)
    bg = min(255, g // 5 + 204)
    bb = min(255, b // 5 + 204)
    return accent, RGBColor(br, bg, bb)


def _rgb_tuple_from_hex(h: str | None, idx: int) -> tuple[int, int, int]:
    raw = (h or "").strip().lstrip("#") or _DEFAULT_ACCENTS[idx % len(_DEFAULT_ACCENTS)].lstrip("#")
    if len(raw) != 6 or not re.fullmatch(r"[0-9a-fA-F]{6}", raw):
        raw = _DEFAULT_ACCENTS[idx % len(_DEFAULT_ACCENTS)].lstrip("#")
    return int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16)


def _body_text_rgb(accent_rgb: tuple[int, int, int]) -> RGBColor:
    """Текст тела: сланцевый оттенок с лёгким подтоном акцента."""
    r, g, b = accent_rgb
    tr = min(255, int(r * 0.18 + 23 * 0.82))
    tg = min(255, int(g * 0.18 + 32 * 0.82))
    tb = min(255, int(b * 0.18 + 48 * 0.82))
    return RGBColor(tr, tg, tb)


# Пресеты: Arial есть в macOS/Windows — Keynote стабильнее, чем с «Calibri Light».
_VISUAL_PRESETS: dict[str, dict[str, Any]] = {
    "corporate": {
        "title_face": "Arial",
        "body_face": "Arial",
        "notes_face": "Arial",
        "title_pt": 32,
        "subtitle_pt": 13,
        "body_pt": 15,
        "footer_pt": 9,
        "bar_in": 1.38,
        "title_bold": False,
    },
    "modern": {
        "title_face": "Arial",
        "body_face": "Arial",
        "notes_face": "Arial",
        "title_pt": 34,
        "subtitle_pt": 14,
        "body_pt": 16,
        "footer_pt": 9,
        "bar_in": 1.42,
        "title_bold": False,
    },
    "bold": {
        "title_face": "Arial",
        "body_face": "Arial",
        "notes_face": "Arial",
        "title_pt": 36,
        "subtitle_pt": 14,
        "body_pt": 16,
        "footer_pt": 10,
        "bar_in": 1.45,
        "title_bold": True,
    },
    "compact": {
        "title_face": "Arial",
        "body_face": "Arial",
        "notes_face": "Arial",
        "title_pt": 28,
        "subtitle_pt": 12,
        "body_pt": 14,
        "footer_pt": 8,
        "bar_in": 1.28,
        "title_bold": False,
    },
}


def _preset_for_slide(slide_data: dict[str, Any], idx: int) -> dict[str, Any]:
    name = str(slide_data.get("visual_style") or slide_data.get("theme") or "").strip().lower()
    base = dict(_VISUAL_PRESETS["corporate"])
    if name in _VISUAL_PRESETS:
        base.update(_VISUAL_PRESETS[name])
    else:
        base.update(_VISUAL_PRESETS[("corporate", "modern", "bold", "compact")[idx % 4]])
    return base


def _set_paragraph_line_spacing(p, multiple: float = 1.12) -> None:
    """Пусто: float line_spacing в python-pptx даёт разметку, с которой Keynote импортёр падает."""
    pass


def parse_slides_json(raw: str) -> list[dict[str, Any]]:
    m = re.search(r"\[.*\]", raw, re.DOTALL)
    if m:
        raw = m.group(0)
    data = json.loads(raw)
    if not isinstance(data, list):
        raise ValueError("slides not a list")
    return [x for x in data if isinstance(x, dict)]


def parse_presentation_json(raw: str) -> tuple[str, list[dict[str, Any]]]:
    """
    Объект {\"deck_title\", \"slides\": [...]} (предпочтительно) или голый массив слайдов.
    """
    t = (raw or "").strip()
    if "```" in t:
        cm = re.search(r"```(?:json)?\s*([\s\S]*?)\s*```", t)
        if cm:
            t = cm.group(1).strip()
    mobj = re.search(r"\{[\s\S]*\"slides\"[\s\S]*\}", t)
    if mobj:
        try:
            obj = json.loads(mobj.group(0))
            slides = obj.get("slides")
            if isinstance(slides, list):
                deck = str(obj.get("deck_title") or obj.get("title") or "").strip()
                return deck, [x for x in slides if isinstance(x, dict)]
        except json.JSONDecodeError:
            pass
    return "", parse_slides_json(raw)


def _is_image_magic(data: bytes) -> bool:
    if len(data) < 12:
        return False
    if data[:8] == b"\x89PNG\r\n\x1a\n":
        return True
    if data[:2] == b"\xff\xd8":
        return True
    if data[:4] == b"GIF8":
        return True
    if data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return True
    return False


async def download_first_web_image(urls: list[str]) -> Optional[Path]:
    """Скачать первое подходящее изображение по URL из поиска."""
    headers = {
        "User-Agent": "Mozilla/5.0 (compatible; GPTHub/1.0; presentation-images)",
        "Accept": "image/avif,image/webp,image/apng,image/*,*/*;q=0.8",
    }
    d = settings.data_dir / "static" / "images"
    d.mkdir(parents=True, exist_ok=True)
    for url in urls[:28]:
        if not (url.startswith("http://") or url.startswith("https://")):
            continue
        try:
            async with httpx.AsyncClient(
                timeout=25.0,
                follow_redirects=True,
                headers=headers,
            ) as client:
                r = await client.get(url)
                r.raise_for_status()
                body = r.content
            if len(body) < 800 or len(body) > 5_000_000:
                continue
            if not _is_image_magic(body):
                continue
            from io import BytesIO

            from PIL import Image

            uid = uuid.uuid4().hex[:12]
            # Keynote часто не принимает WebP/GIF внутри PPTX — сохраняем как PNG/JPEG.
            if body[:2] == b"\xff\xd8":
                p = d / f"webimg_{uid}.jpg"
                p.write_bytes(body)
                return p
            if body[:8] == b"\x89PNG\r\n\x1a\n":
                p = d / f"webimg_{uid}.png"
                p.write_bytes(body)
                return p
            try:
                im = Image.open(BytesIO(body)).convert("RGB")
                p = d / f"webimg_{uid}.png"
                im.save(p, "PNG")
                return p
            except Exception as e:
                logger.debug("skip non-trivial image format: %s", e)
                continue
        except Exception as e:
            logger.debug("skip image url %s: %s", url[:60], e)
            continue
    return None


def write_presentation_sidecar(
    path: Path,
    deck_title: str,
    slides_data: list[dict[str, Any]],
    research_excerpt: str,
    stem: str = "",
) -> None:
    """JSON со структурой для редактора и пересборки PPTX."""
    doc: dict[str, Any] = {
        "deck_title": deck_title,
        "slides": slides_data,
        "research_excerpt": (research_excerpt or "")[:8000],
        "edit_hint": (
            "Веб-редактор: /presentation/editor/?stem=… — после правок нажмите «Пересобрать PPTX». "
            "Либо откройте PPTX в PowerPoint / Keynote."
        ),
    }
    if stem:
        doc["stem"] = stem
        doc["pptx_rel_url"] = f"static/presentations/{stem}.pptx"
        doc["editor_url"] = f"/presentation/editor/?stem={stem}"
    path.write_text(json.dumps(doc, ensure_ascii=False, indent=2), encoding="utf-8")


async def _href_to_local_path(href: str) -> Optional[Path]:
    if not href:
        return None
    if href.startswith("static/"):
        p = settings.data_dir / href
        return p if p.is_file() else None
    if href.startswith("http://") or href.startswith("https://"):
        try:
            async with httpx.AsyncClient(timeout=90.0, follow_redirects=True) as client:
                r = await client.get(href)
                r.raise_for_status()
                body = r.content
                ct = (r.headers.get("content-type") or "").lower()
            if "jpeg" in ct or "jpg" in ct:
                ext = ".jpg"
            elif "webp" in ct:
                ext = ".webp"
            d = settings.data_dir / "static" / "images"
            d.mkdir(parents=True, exist_ok=True)
            fn = f"slide_{uuid.uuid4().hex[:12]}{ext}"
            p = d / fn
            p.write_bytes(body)
            return p
        except Exception as e:
            logger.warning("download slide image url failed: %s", e)
            return None
    return None


def normalize_slide_rows_for_images(slides_data: list[dict[str, Any]], deck_title: str) -> None:
    """Перед поиском: не пустой image_query; кириллица — усилить англ. ключами для DDG."""
    dt = (deck_title or "").strip()
    for row in slides_data:
        if not isinstance(row, dict):
            continue
        iq = (row.get("image_query") or "").strip()
        title = str(row.get("title") or "").strip()
        if len(iq) < 2:
            row["image_query"] = (f"{dt} {title}".strip() if dt else title)[:500]
            iq = (row.get("image_query") or "").strip()
        lat = len(re.findall(r"[a-zA-Z]", iq))
        if _CYRILLIC_RE.search(iq) and lat < 4:
            row["image_query"] = f"{iq} photo stock image"[:500]


def _sanitize_neuro_slide_prompt(prompt: str) -> str:
    """Убрать слова, из-за которых диффузор рисует «слайды» с псевдотекстом."""
    t = _NEURO_SLIDE_JUNK.sub(" ", prompt or "")
    return re.sub(r"\s+", " ", t).strip()[:3500]


def _slide_image_search_queries(row: dict[str, Any], deck_title: str = "") -> list[str]:
    """Несколько запросов для DDG Images: тема + фото/сток/wikimedia."""
    q = (row.get("image_query") or "").strip()
    title = str(row.get("title") or "").strip()
    dt = (deck_title or "").strip()
    base = (q or title)[:500]
    if len(base) < 2 and dt:
        base = f"{dt} {title}".strip()[:500]
    if len(base) < 2:
        return []
    out: list[str] = []
    seen: set[str] = set()

    def add(s: str) -> None:
        s = s.strip()[:500]
        if len(s) < 2:
            return
        key = s.casefold()
        if key in seen:
            return
        seen.add(key)
        out.append(s)

    add(base)
    add(f"{base} photo")
    add(f"{base} stock photo")
    add(f"{base} photography")
    if dt and dt.casefold() not in base.casefold():
        add(f"{dt} {base}"[:500])
    add(f"{base} site:wikimedia.org")
    return out[:10]


async def generate_slide_image(
    client: MWSClient,
    image_model: str,
    prompt_en: str,
) -> Optional[Path]:
    """Фолбэк: одна картинка без текста; не «слайд» и не инфографика."""
    raw = _sanitize_neuro_slide_prompt(prompt_en)
    p = (
        "Single photorealistic photograph or cinematic wide shot, natural lighting, "
        "sharp focus, environmental context, editorial quality. "
        "NOT a slide, NOT a poster, NOT an infographic, NOT a diagram with labels. "
        "ABSOLUTELY NO text, letters, numbers, captions, watermarks, logos, typography, "
        "fake writing, gibberish script, chart labels, UI, speech bubbles. "
        "Pure image only. "
        + raw
        + " . No text anywhere."
    )
    body: dict[str, Any] = {
        "model": image_model,
        "prompt": p[:4000],
        "n": 1,
        "size": "1024x1024",
        "response_format": "b64_json",
        "negative_prompt": (
            "text, words, letters, typography, watermark, logo, caption, title, subtitle, "
            "infographic, chart text, UI, mockup, label, poster, slide, presentation, "
            "gibberish writing, cyrillic text, latin text, numbers on image, signage"
        )[:2000],
    }
    try:
        resp = await client.post_json("/images/generations", body)
        href = await image_api_response_to_sse_href(resp, settings.data_dir)
        return await _href_to_local_path(href)
    except Exception as e:
        logger.warning("slide image generation failed: %s", e)
        if body.pop("negative_prompt", None) is not None:
            try:
                resp = await client.post_json("/images/generations", body)
                href = await image_api_response_to_sse_href(resp, settings.data_dir)
                return await _href_to_local_path(href)
            except Exception as e2:
                logger.warning("slide image generation retry without negative_prompt: %s", e2)
        return None


def _pick_image_model(available_ids: set[str]) -> str:
    mid = settings.image_gen_model
    if mid in available_ids:
        return mid
    for c in ("qwen-image", "qwen-image-lightning", "sd3.5-large-image", "z-image-turbo"):
        if c in available_ids:
            return c
    return mid


def build_colorful_pptx(
    slides_data: list[dict[str, Any]],
    image_paths: list[Optional[Path]],
    out_path: Path,
    deck_title: str = "",
) -> None:
    """Собрать PPTX: типографика, цвета из акцента, полосы, подвал, заметки."""
    prs = _open_presentation_base()
    try:
        blank = prs.slide_layouts[6]
    except IndexError:
        blank = prs.slide_layouts[-1]

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    n_slides = len(slides_data)

    for idx, slide_data in enumerate(slides_data):
        preset = _effective_preset(slide_data, idx)
        bar_h = Inches(float(preset["bar_in"]))
        raw_img = image_paths[idx] if idx < len(image_paths) else None
        img_path = _ensure_keynote_safe_image(raw_img)
        title = sanitize_ooxml_text(str(slide_data.get("title") or "Слайд"), 500)
        bullets = slide_data.get("bullets") or []
        if not isinstance(bullets, list):
            bullets = []
        accent_s = slide_data.get("accent")
        if isinstance(accent_s, str):
            accent_s = accent_s.strip()
        else:
            accent_s = None

        accent, light_bg = _hex_to_rgb(accent_s, idx)
        accent_rgb = _rgb_tuple_from_hex(accent_s, idx)
        body_col = _body_text_rgb(accent_rgb)

        slide = prs.slides.add_slide(blank)

        # Подложка контента
        body_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
        body_rect.fill.solid()
        body_rect.fill.fore_color.rgb = light_bg
        body_rect.line.fill.background()

        # Левый акцентный штрих (как у корпоративных шаблонов)
        strip_w = Inches(0.07)
        left_strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            bar_h,
            strip_w,
            Emu(int(slide_h) - int(bar_h)),
        )
        left_strip.fill.solid()
        left_strip.fill.fore_color.rgb = accent
        left_strip.line.fill.background()

        # Верхняя полоса
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, bar_h)
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = accent
        top_bar.line.fill.background()

        # Заголовок
        tx_title = slide.shapes.add_textbox(
            Inches(0.42),
            Inches(0.16),
            Inches(9.2),
            Inches(1.05),
        )
        tf_t = tx_title.text_frame
        tf_t.word_wrap = True
        tf_t.margin_bottom = Pt(2)
        p0 = tf_t.paragraphs[0]
        p0.text = title
        p0.font.name = preset["title_face"]
        p0.font.size = Pt(int(preset["title_pt"]))
        p0.font.bold = bool(preset.get("title_bold", False))
        p0.font.color.rgb = RGBColor(255, 255, 255)
        p0.alignment = PP_ALIGN.LEFT
        _set_paragraph_line_spacing(p0, 1.05)

        subtitle = slide_data.get("subtitle")
        if isinstance(subtitle, str) and subtitle.strip():
            ps = tf_t.add_paragraph()
            ps.text = sanitize_ooxml_text(subtitle.strip(), 400)
            ps.font.name = preset["body_face"]
            ps.font.size = Pt(int(preset["subtitle_pt"]))
            ps.font.bold = False
            ps.font.italic = False
            ps.font.color.rgb = RGBColor(235, 240, 255)
            ps.space_before = Pt(3)
            _set_paragraph_line_spacing(ps, 1.1)

        left_margin = Inches(0.52)
        top_body = bar_h + Inches(0.28)
        body_h = Inches(5.75)
        if img_path and img_path.is_file():
            text_w = Inches(4.75)
            tx_body = slide.shapes.add_textbox(left_margin, top_body, text_w, body_h)
        else:
            tx_body = slide.shapes.add_textbox(left_margin, top_body, Inches(9.0), body_h)

        tf_b = tx_body.text_frame
        tf_b.word_wrap = True
        tf_b.vertical_anchor = MSO_ANCHOR.TOP
        tf_b.margin_left = Inches(0.06)
        tf_b.margin_right = Inches(0.06)
        tf_b.margin_top = Pt(2)

        first_line = True
        for bullet in bullets[:12]:
            text = sanitize_ooxml_text(str(bullet).strip(), 2000)
            if not text:
                continue
            if first_line:
                p = tf_b.paragraphs[0]
                first_line = False
            else:
                p = tf_b.add_paragraph()
            # Одна строка с маркером — так Keynote не ломается на пустом p.text + нескольких run.
            p.text = "- " + text
            p.font.name = preset["body_face"]
            p.font.size = Pt(int(preset["body_pt"]))
            p.font.color.rgb = body_col
            p.space_after = Pt(10)
            p.level = 0
            _set_paragraph_line_spacing(p, 1.14)

        if not bullets:
            p = tf_b.paragraphs[0]
            p.text = "—"
            p.font.name = preset["body_face"]
            p.font.size = Pt(int(preset["body_pt"]))
            p.font.color.rgb = RGBColor(130, 130, 140)

        if img_path and img_path.is_file():
            pic_left = Inches(5.42)
            pic_top = top_body
            pic_w = Inches(4.2)
            try:
                slide.shapes.add_picture(str(img_path), pic_left, pic_top, width=pic_w, height=body_h)
            except Exception as e:
                logger.warning("add_picture failed, slide without image: %s", e)

        # Подвал: название деки + номер слайда
        foot_top = Emu(int(slide_h) - int(Inches(0.4)))
        if (deck_title or "").strip():
            ft_left = slide.shapes.add_textbox(
                Inches(0.42),
                foot_top,
                Inches(7.5),
                Inches(0.32),
            )
            tff = ft_left.text_frame
            fp = tff.paragraphs[0]
            fp.text = sanitize_ooxml_text((deck_title or "").strip(), 140)
            fp.font.name = preset["body_face"]
            fp.font.size = Pt(int(preset["footer_pt"]))
            fp.font.color.rgb = RGBColor(110, 118, 128)
            fp.alignment = PP_ALIGN.LEFT
        fn_box = slide.shapes.add_textbox(
            Inches(8.85),
            foot_top,
            Inches(1.25),
            Inches(0.32),
        )
        fnp = fn_box.text_frame.paragraphs[0]
        fnp.text = sanitize_ooxml_text(f"{idx + 1} / {n_slides}", 32)
        fnp.font.name = preset["body_face"]
        fnp.font.size = Pt(int(preset["footer_pt"]))
        fnp.font.color.rgb = RGBColor(150, 155, 165)
        fnp.alignment = PP_ALIGN.RIGHT

        sn_parts: list[str] = []
        sn = slide_data.get("speaker_notes") or slide_data.get("notes")
        if isinstance(sn, str) and sn.strip():
            sn_parts.append(sanitize_ooxml_text(sn.strip(), 12000))
        sources = slide_data.get("sources")
        if isinstance(sources, list):
            lines: list[str] = []
            for s in sources[:8]:
                if isinstance(s, dict):
                    t = sanitize_ooxml_text(str(s.get("title") or "").strip(), 500)
                    u = sanitize_ooxml_text(str(s.get("url") or "").strip(), 2000)
                    if u:
                        lines.append(f"• {t}: {u}" if t else f"• {u}")
                elif isinstance(s, str) and s.strip():
                    lines.append("• " + sanitize_ooxml_text(s.strip(), 2000))
            if lines:
                sn_parts.append("Источники:\n" + "\n".join(lines))
        if sn_parts:
            try:
                ns = slide.notes_slide
                # Одна вертикальная «простыня» с \n — Keynote хуже переваривает несколько абзацев из \n\n.
                note_blob = "\n".join(sn_parts)
                note_blob = re.sub(r"\n{3,}", "\n\n", note_blob)
                ns.notes_text_frame.text = sanitize_ooxml_text(note_blob, 10000)
                for np in ns.notes_text_frame.paragraphs:
                    np.font.name = preset["notes_face"]
                    np.font.size = Pt(11)
                    np.font.color.rgb = RGBColor(55, 55, 62)
                    _set_paragraph_line_spacing(np, 1.12)
            except Exception as e:
                logger.warning("speaker notes: %s", e)

    prs.save(str(out_path))
    n = len(slides_data)
    if settings.gena_pptx_validate_zip and not _validate_ooxml_package(out_path, n):
        logger.warning("pptx package validation failed: %s", out_path)
    if settings.gena_pptx_roundtrip:
        try:
            _roundtrip_repair(out_path)
        except Exception as e:
            logger.warning("pptx round-trip repair failed: %s", e)
        if settings.gena_pptx_validate_zip and not _validate_ooxml_package(out_path, n):
            logger.warning("pptx package invalid after round-trip: %s", out_path)


async def _resolve_one_slide_image(
    client: MWSClient,
    row: dict[str, Any],
    model_id: str,
    deck_title: str = "",
) -> Optional[Path]:
    """Картинка: сначала веб (DDG), затем нейро. Режим влияет только на обязательность нейро-fallback."""
    mode = str(row.get("image_mode") or "auto").strip().lower()

    async def gen_neuro() -> Optional[Path]:
        ip = row.get("image_prompt")
        title = str(row.get("title") or "")
        if isinstance(ip, str) and ip.strip():
            prompt = _sanitize_neuro_slide_prompt(ip.strip())
        else:
            prompt = _sanitize_neuro_slide_prompt(
                f"Photograph of subject matter: {title}. Context: {row.get('bullets', [])}"
            )
        try:
            return await generate_slide_image(client, model_id, prompt)
        except Exception as e:
            logger.warning("slide neuro image: %s", e)
            return None

    async def try_web() -> Optional[Path]:
        """Несколько формулировок запроса — иначе DDG часто отдаёт пусто или битые URL."""
        for sq in _slide_image_search_queries(row, deck_title):
            urls = image_search_ddg_urls(sq, max_results=28)
            if not urls:
                logger.debug("slide image DDG: no urls for query=%r", sq[:80])
                continue
            got = await download_first_web_image(urls)
            if got:
                logger.info("slide image: web ok query=%r", sq[:100])
                return got
        logger.warning("slide image: web search failed for all queries, falling back to neuro")
        return None

    # Всегда в приоритете реальные изображения из сети (в т.ч. при image_mode=generate —
    # нейро только если веб не дал файла).
    web_got = await try_web()
    if web_got:
        return web_got

    if mode in ("search", "web", "internet", "ddg"):
        return await gen_neuro()

    if mode in ("generate", "ai", "neuro", "neural"):
        return await gen_neuro()

    # auto и прочее: как и раньше — нейро после неудачного веба
    return await gen_neuro()


async def resolve_slide_images_progress(
    client: MWSClient,
    slides_data: list[dict[str, Any]],
    available_ids: set[str],
    deck_title: str = "",
):
    """По мере готовности каждого изображения: (индекс слайда, путь или None). Параллельно до 3 потоков."""
    if not slides_data:
        return
    model_id = _pick_image_model(available_ids)
    sem = asyncio.Semaphore(3)

    async def wrapped(i: int, row: dict[str, Any]) -> tuple[int, Optional[Path]]:
        async with sem:
            p = await _resolve_one_slide_image(client, row, model_id, deck_title=deck_title)
        return i, p

    coros = [wrapped(i, row) for i, row in enumerate(slides_data)]
    for coro in asyncio.as_completed(coros):
        i, p = await coro
        yield i, p


async def resolve_slide_images(
    client: MWSClient,
    slides_data: list[dict[str, Any]],
    available_ids: set[str],
    deck_title: str = "",
) -> list[Optional[Path]]:
    """По одному изображению на слайд: веб и/или генерация."""
    n = len(slides_data)
    paths: list[Optional[Path]] = [None] * n
    async for i, p in resolve_slide_images_progress(
        client, slides_data, available_ids, deck_title=deck_title
    ):
        paths[i] = p
    return paths
