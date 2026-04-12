"""
Фичи gena/router: стрим презентаций PPTX, стрим картинок, стрим deep research, SSE-хелперы.
"""

from __future__ import annotations

import json
import logging
import re
import uuid
from typing import AsyncGenerator

import httpx
from fastapi import Request

from app.config import settings
from app.image_utils import image_api_response_to_sse_href
from app.mws_client import MWSClient
from app.router_logic import IMAGE_GEN_RE, MUSIC_GEN_RE, PRESENTATION_RE, gena_chat_target
from app.web_tools import (
    deep_research_ddg,
    extract_urls,
    fetch_url_text,
    should_run_deep_research,
)

logger = logging.getLogger("gpthub.gena")


def sse_delta(content: str) -> str:
    esc = json.dumps(content, ensure_ascii=False)
    return f'data: {{"choices": [{{"delta": {{"content": {esc} }}}}]}}\n\n'


def _pick_model(preferred: str, available: set[str], fallback: str) -> str:
    if preferred in available:
        return preferred
    if fallback in available:
        return fallback
    for x in sorted(available):
        if x != settings.auto_model_id:
            return x
    return fallback


def public_static_url(request: Request, rel_path: str) -> str:
    """URL для скачивания файлов из /static/…"""
    base = (settings.public_base_url or "").strip().rstrip("/")
    if base:
        return f"{base}/{rel_path.lstrip('/')}"
    return str(request.base_url).rstrip("/") + "/" + rel_path.lstrip("/")


async def stream_presentation_pptx(
    request: Request,
    client: MWSClient,
    prompt: str,
    available_ids: set[str],
) -> AsyncGenerator[str, None]:
    from pptx import Presentation

    yield sse_delta("*(Презентация: генерирую структуру слайдов…)*\n\n")
    model = _pick_model(settings.gena_code_model, available_ids, settings.default_llm)
    system_prompt = (
        "Ты — генератор структуры презентаций. Верни СТРОГО JSON-массив объектов. "
        "Каждый объект: 'title' (строка), 'bullets' (массив строк). "
        "Не менее 5 слайдов. Только JSON, без markdown."
    )
    try:
        data = await client.post_json(
            "/chat/completions",
            {
                "model": model,
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt[:8000]},
                ],
                "temperature": 0.3,
                "max_tokens": 4000,
            },
        )
        raw = (data.get("choices") or [{}])[0].get("message", {}).get("content") or ""
        m = re.search(r"\[.*\]", raw, re.DOTALL)
        if m:
            raw = m.group(0)
        slides_data = json.loads(raw)
        if not isinstance(slides_data, list):
            raise ValueError("slides not a list")

        yield sse_delta("*(Презентация: собираю PPTX…)*\n\n")

        static_dir = settings.data_dir / "static" / "presentations"
        static_dir.mkdir(parents=True, exist_ok=True)
        fname = f"presentation_{uuid.uuid4().hex[:10]}.pptx"
        fpath = static_dir / fname

        prs = Presentation()
        for slide_data in slides_data:
            if not isinstance(slide_data, dict):
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = str(slide_data.get("title", "Слайд"))
            body = slide.placeholders[1]
            tf = body.text_frame
            bullets = slide_data.get("bullets") or []
            if isinstance(bullets, list):
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        tf.text = str(bullet)
                    else:
                        p = tf.add_paragraph()
                        p.text = str(bullet)

        prs.save(str(fpath))
        url = public_static_url(request, f"static/presentations/{fname}")
        yield sse_delta(
            f"✅ **Презентация готова.**\n\n[Скачать PPTX]({url})\n\n"
        )
    except Exception as e:
        logger.exception("presentation")
        yield sse_delta(f"Ошибка презентации: {e}\n\n")
    yield "data: [DONE]\n\n"


async def stream_music_demo(
    request: Request,
    client: MWSClient,
    prompt: str,
    available_ids: set[str],
) -> AsyncGenerator[str, None]:
    """SSE: статусы + ссылка на демо MP3 (как у картинки, но для музыки)."""
    from app.music_demo import build_mp3_from_prompt, melody_notes_from_llm

    yield sse_delta("*(Демо-мелодия: подбираю ноты…)*\n\n")
    mid = gena_chat_target()
    if mid not in available_ids:
        if settings.default_llm in available_ids:
            mid = settings.default_llm
        else:
            mid = next(iter(sorted(available_ids - {settings.auto_model_id})), settings.default_llm)
    try:
        llm_notes = await melody_notes_from_llm(client, prompt, mid)
        yield sse_delta("*(Демо-мелодия: синтез MP3…)*\n\n")
        mp3 = build_mp3_from_prompt(prompt, llm_notes)
    except Exception as e:
        logger.exception("music demo stream")
        yield sse_delta(f"Не удалось сгенерировать MP3: {e}\n\n")
        yield "data: [DONE]\n\n"
        return

    static_dir = settings.data_dir / "static" / "music"
    static_dir.mkdir(parents=True, exist_ok=True)
    fname = f"demo_{uuid.uuid4().hex[:12]}.mp3"
    (static_dir / fname).write_bytes(mp3)
    url = public_static_url(request, f"static/music/{fname}")
    yield sse_delta(
        "Демо-мелодия (простой синтез по нотам, не студийный саундтрек):\n\n"
        f"[Скачать MP3]({url})\n\n"
    )
    yield "data: [DONE]\n\n"


async def stream_image_markdown(
    request: Request,
    client: MWSClient,
    prompt: str,
    available_ids: set[str],
) -> AsyncGenerator[str, None]:
    yield sse_delta("*(Генерация изображения…)*\n\n")
    model_id = settings.image_gen_model
    if model_id not in available_ids:
        for c in ("qwen-image", "qwen-image-lightning", "sd3.5-large-image", "z-image-turbo"):
            if c in available_ids:
                model_id = c
                break
    try:
        enhance = await client.post_json(
            "/chat/completions",
            {
                "model": _pick_model(gena_chat_target(), available_ids, settings.default_llm),
                "messages": [
                    {
                        "role": "system",
                        "content": "Output ONLY a concise English image generation prompt, no other text.",
                    },
                    {"role": "user", "content": prompt[:2000]},
                ],
                "max_tokens": 500,
            },
        )
        ep = (
            (enhance.get("choices") or [{}])[0]
            .get("message", {})
            .get("content", "")
            .strip()
        )
        if ep:
            prompt = ep
    except Exception:
        pass

    try:
        img_resp = await client.post_json(
            "/images/generations",
            {
                "model": model_id,
                "prompt": prompt[:4000],
                "n": 1,
                "size": "1024x1024",
                "response_format": "b64_json",  # сразу base64 → сохраняем в static, не в SSE
            },
        )
        href = await image_api_response_to_sse_href(img_resp, settings.data_dir)
        if href.startswith("http://") or href.startswith("https://"):
            display = href
        elif href.startswith("static/"):
            display = public_static_url(request, href)
        else:
            display = href
        if display:
            yield sse_delta(f"![Изображение]({display})\n\n")
        else:
            yield sse_delta("Не удалось получить ссылку на изображение.\n\n")
    except Exception as e:
        yield sse_delta(f"Ошибка генерации: {e}\n\n")
    yield "data: [DONE]\n\n"


async def stream_deep_research(
    client: MWSClient,
    user_prompt: str,
    available_ids: set[str],
) -> AsyncGenerator[str, None]:
    yield sse_delta("*(Deep Research: собираю источники…)*\n\n")
    block = deep_research_ddg(user_prompt)
    urls = extract_urls(block)[:3]
    fetched: list[str] = []
    for u in urls:
        t = await fetch_url_text(u, max_chars=6000)
        fetched.append(f"=== {u} ===\n{t}")
    ctx = block + "\n\n" + "\n\n".join(fetched)
    yield sse_delta("*(Deep Research: пишу отчёт…)*\n\n")

    model = _pick_model(settings.gena_long_doc_model, available_ids, settings.default_llm)
    sys_msg = (
        "Ты — исследователь. По теме пользователя и контексту из веба дай структурированный отчёт в Markdown.\n\n"
        f"КОНТЕКСТ:\n{ctx[:24000]}"
    )
    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": sys_msg},
            {"role": "user", "content": user_prompt[:8000]},
        ],
        "stream": True,
        "temperature": 0.4,
    }
    headers = {
        "Authorization": f"Bearer {settings.mws_api_key}",
        "Content-Type": "application/json",
    }
    done = False
    async with httpx.AsyncClient(timeout=300.0) as http:
        async with http.stream(
            "POST",
            f"{settings.mws_api_base.rstrip('/')}/chat/completions",
            headers=headers,
            content=json.dumps(payload),
        ) as resp:
            if resp.status_code >= 400:
                err = await resp.aread()
                yield sse_delta(f"Ошибка API: {err.decode()[:500]}\n\n")
                yield "data: [DONE]\n\n"
                return
            async for line in resp.aiter_lines():
                if not line:
                    continue
                if line.startswith("data:"):
                    pl = line[5:].lstrip()
                    if pl == "[DONE]":
                        done = True
                        yield "data: [DONE]\n\n"
                        break
                    yield line + "\n\n"
            if not done:
                yield "data: [DONE]\n\n"


def should_stream_presentation(last_text: str, stream: bool) -> bool:
    return bool(stream and last_text and PRESENTATION_RE.search(last_text))


def should_stream_deep_gena(last_text: str, stream: bool) -> bool:
    return bool(stream and last_text and should_run_deep_research(last_text))


def should_stream_music_gena(
    last_text: str, stream: bool, has_image: bool, has_audio: bool
) -> bool:
    return bool(
        stream
        and last_text
        and not has_image
        and not has_audio
        and MUSIC_GEN_RE.search(last_text)
    )


def should_stream_image_gena(last_text: str, stream: bool, has_image: bool) -> bool:
    return bool(
        stream
        and last_text
        and not has_image
        and IMAGE_GEN_RE.search(last_text)
        and not MUSIC_GEN_RE.search(last_text)
    )
